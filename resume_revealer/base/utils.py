
import fitz
import nltk
import re
import os
from pptx import Presentation  
import docx  
from PIL import Image
import pytesseract
import aspose.slides as slides
import aspose.words as words



def extract_text_from_pdf(pdf_path):
    text_content = ''
    try:
        with fitz.open(pdf_path) as pdf_document:
            num_pages = pdf_document.page_count
            for page_num in range(num_pages):
                page = pdf_document.load_page(page_num)
                text_content += page.get_text()
    except FileNotFoundError:
        print(f"Error: File '{pdf_path}' not found.")
    except Exception as e:
        print("An error occurred:", str(e))

    # return text_content.split('\n')
    lines = text_content.split('\n')
    lines_and_commas = [item.strip() for line in lines for item in line.split(',')]
    return lines_and_commas

def extract_text_from_ppt(ppt_path):
    text_content = ''
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text
    except FileNotFoundError:
        print(f"Error: File '{ppt_path}' not found.")
    except Exception as e:
        print("An error occurred:", str(e))
    # return text_content.split('\n')
    return text_content

def extract_text_from_doc(doc_path):
    text_content = ''
    try:
        doc = docx.Document(doc_path)
        for paragraph in doc.paragraphs:
            text_content += paragraph.text
    except FileNotFoundError:
        print(f"Error: File '{doc_path}' not found.")
    except Exception as e:
        print("An error occurred:", str(e))
    # return text_content.split('\n')
    lines = text_content.split('\n')
    lines_and_commas = [item.strip() for line in lines for item in line.split(',')]
    return lines_and_commas

def extract_text_from_jpg(jpg_path):
    text_content = ''
    try:
        with Image.open(jpg_path) as img:
            text_content = pytesseract.image_to_string(img)
    except FileNotFoundError:
        print(f"Error: File '{jpg_path}' not found.")
    except Exception as e:
        print("An error occurred:", str(e))
    return text_content.split('\n')

# def preprocess_text(text):
#     # stop = set(stopwords.words('english'))
#     # text = re.sub(r'[^A-Za-z\s\']', '', text)  # Remove non-alphabetic characters except spaces and apostrophes
#     # text = nltk.word_tokenize(text)  # Tokenize the text
#     # text = [word.replace("'", "") for word in text]  # Remove apostrophes
#     # text = [word for word in text if word.lower() not in stop]  # Remove stopwords
#     # return text
#     text = BeautifulSoup(text, "html.parser").get_text()
    
#     # Remove non-alphabetic characters except spaces and apostrophes
#     text = re.sub(r'[^A-Za-z\s\']', '', text)
    
#     # Tokenize the text
#     text = nltk.word_tokenize(text)
    
#     # Remove apostrophes
#     text = [word.replace("'", "") for word in text]
    
#     # Remove stopwords
#     stop = set(stopwords.words('english'))
#     text = [word for word in text if word.lower() not in stop]
    
#     return text

# def nlp_analysis(text):
#     # Perform NLP analysis here (e.g., Named Entity Recognition, sentiment analysis, etc.)
#     # Return the results
#     # Example:
#     nlp_results = {
#         'skills': ['Python', 'Machine Learning', 'Data Analysis'],
#         'experience': '5 years',
#         'education': 'Bachelor\'s degree in Computer Science',
#     }
#     return nlp_results